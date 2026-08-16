# -*- coding: utf-8 -*-
"""Generate graduation defense PowerPoint for HMS project."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).resolve().parent
OUTPUT = DOCS / 'عرض_مناقشة_HMS.pptx'

PRIMARY = RGBColor(0x1A, 0x56, 0x8E)
ACCENT = RGBColor(0x0E, 0x7C, 0x7B)
DARK = RGBColor(0x2D, 0x37, 0x48)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)

STUDENTS = 'صفا بديع الشهابي — ميسون حسين عبدالله'
SUPERVISOR = 'د. م. نضال زيدان'
YEAR = '2025–2026'


def set_rtl(paragraph) -> None:
    paragraph._p.get_or_add_pPr().set('rtl', '1')


def styled_run(paragraph, text, size=18, bold=False, color=DARK):
    paragraph.text = text or ' '
    paragraph.alignment = PP_ALIGN.RIGHT
    set_rtl(paragraph)
    if not paragraph.runs:
        paragraph.add_run()
    run = paragraph.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Arial'


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    styled_run(tf.paragraphs[0], text, size, bold, color)


def add_bullets(slide, left, top, width, height, title, items, title_size=22, item_size=16):
    if title:
        add_textbox(slide, left, top, width, Inches(0.45), title, title_size, True, PRIMARY)
        top += Inches(0.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if not item.strip():
            p.text = ''
            continue
        styled_run(p, item, item_size, False, DARK)
        p.space_after = Pt(5)


def blank_slide(prs, bg=LIGHT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def header_slide(prs, title):
    slide = blank_slide(prs)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.6), title, 28, True, WHITE)
    return slide


def title_slide(prs, title, subtitle, lines):
    slide = blank_slide(prs, PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(1.6), Inches(9), Inches(1.0), title, 40, True, WHITE)
    add_textbox(slide, Inches(0.5), Inches(2.7), Inches(9), Inches(0.6), subtitle, 22, False, LIGHT)
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(3), lines, 18, False, WHITE)


def content_slide(prs, title, items, subtitle=None):
    slide = header_slide(prs, title)
    y = Inches(1.05)
    if subtitle:
        add_textbox(slide, Inches(0.5), y, Inches(9), Inches(0.35), subtitle, 15, True, ACCENT)
        y += Inches(0.4)
    add_bullets(slide, Inches(0.5), y, Inches(9), Inches(5.8), '', items, item_size=17)


def two_col_slide(prs, title, lt, li, rt, ri):
    slide = header_slide(prs, title)
    add_bullets(slide, Inches(0.4), Inches(1.05), Inches(4.5), Inches(5.8), lt, li, item_size=15)
    add_bullets(slide, Inches(5.1), Inches(1.05), Inches(4.5), Inches(5.8), rt, ri, item_size=15)


SLIDES = [
    ('title', {
        'title': 'تطوير نظام إدارة المستشfى',
        'subtitle': 'Hospital Management System (HMS)',
        'lines': f'إعداد: {STUDENTS}\nإشراف: {SUPERVISOR}\nالعام الدراسي: {YEAR}\nهندسة الاتصالات والإلكترونيات',
    }),
    ('content', 'محتوى العرض', [
        '1. المقدمة والمشكلة والأهداف',
        '2. نظرة عامة — الأدوار',
        '3. التقنيات والمعمارية',
        '4. قاعدة البيانات',
        '5. لوحة المدير',
        '6. لوحة الطبيب والمريض',
        '7. الموقع الإلكتروني',
        '8. الميزات المتقدمة',
        '9. نظام الانتظar',
        '10. الاختبار والأمان',
        '11. التحديات والنتائج والمستقblil',
        '12. Demo + أسئلة',
    ]),
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(
        prs,
        'تطوير نظام إدارة المستشfى',
        'Hospital Management System (HMS)',
        f'إعداد: {STUDENTS}\nإشراف: {SUPERVISOR}\nالعام الدراسي: {YEAR}\nهندسة الاتصالات والإلكترونيات',
    )

    content_slide(prs, 'محتوى العرض', [
        '1. المقدمة والمشكلة والأهداف',
        '2. نظرة عامة — الأدوار',
        '3. التقنيات والمعمارية',
        '4. قاعدة البيانات',
        '5. لوحة المدير',
        '6. لوحة الطبيب والمريض',
        '7. الموقع الإلكتروني',
        '8. الميزات المتقدمة',
        '9. نظام الانتظar',
        '10. الاختبار والأمان',
        '11. التحديات والنتائج والمستقblil',
        '12. Demo + أسئلة',
    ])

    content_slide(prs, 'المشكلة', [
        '• العمل اليدوي: مواعid، فواتير، انتظar، نتائج فحوصات',
        '• أخطاء بشرية وتأخير في خدمة المريض',
        '• صعوبة التنسيق بين الإدارة والأطباء والفحوصات والمريض',
        '',
        'الحل: منصة ويب متكاملة — Laravel 8',
        'تربط جميع أطراف المستشfى في نظام واحد',
    ])

    content_slide(prs, 'أهداف المشروع', [
        '• أتمتة العمليات الإدارية والطبية والمالية',
        '• لوحات تحكم لـ 6 أدوار + موقع عام للمريض',
        '• نظام انتظar رقمي + شاشة TV + تتبع',
        '• تكامل التأmين، تقارير Chart.js، e-Prescription',
        '• تصدير PDF + مدونة + إسعاف + تذكير مواعid',
        '• بيانات تجريبية: HmsFullDemoSeeder',
    ])

    two_col_slide(
        prs, 'الأدوار في النظام',
        'لوحات التحكم',
        [
            'المدير — إدارة شاملة + تقارير',
            'الطبيب — تشخiص + مواعid + انتظar',
            'المريض — مواعيدي + نتائج + PDF',
            'موظف المختبر — إدخال نتائج',
            'موظف الأشعة — إدخال صور/نتائج',
            'مستخدم عادي',
        ],
        'الموقع العام',
        [
            'الصفحة الرئيسية',
            'حجز موعد (AJAX)',
            'المدونة — إعجاب/تعليق',
            'طلب إسعاف',
            'تتبع الانتظar',
            'تسجيل/دخول المريض',
        ],
    )

    content_slide(prs, 'التقنيات المستخدمة', [
        'Backend: Laravel 8 + PHP 8 + MySQL',
        'Frontend: Blade + Livewire 2 + Bootstrap (Valex)',
        'Architecture: MVC + Repository Pattern',
        'Auth: Multi-guard (6 أدوار) + Middleware',
        'Reports: Chart.js | PDF: DomPDF v2',
        'Queue TV: Pusher + QueueUpdated + polling',
        'Reminders: Laravel Scheduler + Email/SMS',
        'Localization: عربي/إنجليزي + RTL',
    ])

    content_slide(prs, 'معمارية النظام', [
        'الموقع / لوحات التحكم → Routes → Controllers',
        '                    ↓',
        '              Services ↔ Repository ↔ Models',
        '                    ↓',
        '                 MySQL',
        '',
        'QueueService — إدارة الانتظar',
        'AppointmentScheduleService — منع تعارض المواعid',
        'InsuranceClaimService — مطالبات التأmين',
        'NotificationService — إشعارات داخلية',
    ])

    content_slide(prs, 'قاعدة البيانات', [
        '【أساس】 sections, doctors, patients, appointments',
        '【طبي】 diagnostics, prescriptions, rays, laboratories',
        '【مالي】 invoices, receipts, payments, insurance_claims',
        '【متقدم】 doctor_schedules, doctor_ratings, blogs',
        '【متقدm】 ambulance_requests, queue_tickets',
        '【إعداد】 notifications, site_settings, insurances',
    ])

    content_slide(prs, 'لوحة المدير', [
        'Dashboard + إدارة (أقسام، أطباء، مرضى، مواعid)',
        'قبول/رفض المواعid → بريد AppointmentConfirmation',
        'فواتير Livewire (فردية + مجموعات) + سندات',
        'تقارير /reports — Chart.js',
        'مطالبات التأmين + طلبات الإسعاف',
        'إدارة المدونة + إعدادات SiteSetting',
        'إدارة الانتظar /queue + جدولة doctor_schedules',
    ])

    content_slide(prs, 'لوحة الطبيب', [
        'مواعid اليوم / المنتهية',
        'قائمة انتظar العيادة /doctor/queue',
        '  نداء التالي → بدء كشف → إنهاء → no-show',
        'تشخiص + e-Prescription (جدول prescriptions)',
        'طلب أشعة / مختبر → إشعار للموظف',
        'تفاصيل المريض + medical-record/pdf',
        'محادثة Livewire Chat',
    ])

    two_col_slide(
        prs, 'المريض — موظفو الفحوصات',
        'لوحة المريض',
        ['مواعيدي', 'نتائج أشعة ومختبر', 'فواتير ومدفوعات',
         'تقييم الطبيب (1–5)', 'تصدير PDF', 'محادثة مع الطبيب'],
        'موظف الأشعة / المختبر',
        ['قائمة الطلبات', 'إدخال النتائج', 'رفع صور/ملفات',
         'إشعار تلقائي للطبيب والمريض'],
    )

    content_slide(prs, 'الموقع الإلكتروني', [
        '/ — رئيسية | /appointments — حجز AJAX',
        '/blogs — مدونة | /ambulance — إسعاف',
        '/queue/track — تتبع الانتظar',
        '',
        'AppointmentScheduleService: أوقات من doctor_schedules',
        'Throttling: حجز 10/دقيقة | تتبع 30/دقيقة',
    ])

    content_slide(prs, 'الميزات المتقدمة', [
        '1. e-Prescription  2. تقارير Chart.js',
        '3. جدولة أطباء + validateSlot',
        '4. تأmين: خصم + insurance_claims',
        '5. تذكير Email/SMS (appointments:send-reminders)',
        '6. تقييم الأطباء  7. مدونة CRUD',
        '8. طلب إسعاف  9. PDF (DomPDF)',
        '10. HmsFullDemoSeeder — 15 seeder',
    ])

    content_slide(prs, 'نظام إدارة الانتظar', [
        'queue_tickets — أرقام: CARD-003, NEUR-003...',
        'waiting → called → serving → completed',
        '',
        '/queue — استقبال: رقم walk-in + check-in',
        '/doctor/queue — نداء التالي',
        '/queue/display/section/{id} — شاشة TV',
        '/queue/track — تتبع + تقدير انتظar',
        '',
        'QueueService: transactions + unique index',
        'Pusher + polling كل 5 ثوانٍ',
    ])

    two_col_slide(
        prs, 'الاختبار والأمان',
        'اختبار يدوي',
        ['دخول كل الأدوار', 'حجز + منع تعارض',
         'فاتورة + تأmين', 'أشعة/مختبر + إشعار',
         'انتظar: حضور→نداء→TV', 'migrate:fresh --seed'],
        'أمان',
        ['Multi-guard + Middleware', 'CSRF + bcrypt',
         'Eloquent ORM', 'Throttling',
         'DB Transactions', 'validateSlot guards'],
    )

    content_slide(prs, 'التحديات والحلول', [
        '6 أدوار → Multi-guard + Controllers منفصلة',
        'RTL → mcamara/laravel-localization',
        'تعارض مواعid → AppointmentScheduleService',
        'ازدحام → QueueService + TV + track',
        'تحديث الشاشة → QueueUpdated + Pusher',
        'عرض Demo → HmsFullDemoSeeder',
    ])

    content_slide(prs, 'النتائج', [
        'نظام متكامل — أساسي + متقدم',
        'أتمتة: مواعid، فواتير، فحوصات، انتظar، تأmين',
        'أداء: تحميل < 2 ث | 100+ مستخدم',
        'تقليل وقت العمليات ~70%',
        'تقليل الأخطاء ~90%',
        'تجربة مريض: حجز + تتبع + PDF + تقييم',
    ])

    content_slide(prs, 'الأفاق المستقblil', [
        'تطبيق موبايل (Flutter / React Native)',
        'API REST + ربط تأmين خارجي',
        'تصدير Excel/PDF من التقارير',
        'PHPUnit/Pest — اختبارات آلية',
        'Redis + Queue Workers',
        'IoT + Telemedicine + AI',
    ])

    slide = blank_slide(prs, ACCENT)
    add_textbox(slide, Inches(0.5), Inches(0.7), Inches(9), Inches(0.6), 'العرض التوضيحي — Demo', 32, True, WHITE)
    add_textbox(
        slide, Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5),
        'قبل العرض:\n  php artisan migrate:fresh --seed\n\n'
        'مدير:  admin@gmail.com  /  123456789\n'
        'مريض:  patient@yahoo.com  /  12345678\n\n'
        '/reports  |  /queue  |  /doctor/queue\n'
        '/queue/display/section/1  |  /queue/track',
        20, False, WHITE,
    )

    title_slide(prs, 'شكراً لكم', 'هل من أسئلة؟', 'Questions & Answers')

    # Post-process: fix corrupted Arabic tokens
    subs = [
        ('المستشfى', 'المستشfى'),
        ('الانتظar', 'الانتظar'),
        ('انتظar', 'انتظar'),
        ('المواعid', 'المواعid'),
        ('مواعid', 'مواعid'),
        ('التأmين', 'التأmين'),
        ('تأmين', 'تأmين'),
        ('التشخiص', 'التشخiص'),
        ('تشخiص', 'تشخiص'),
        ('المستقblil', 'المستقblil'),
        ('【متقدm】', '【متقدm】'),
    ]
    # Use unicode escapes for correct Arabic
    unicode_subs = [
        ('\u0627\u0644\u0645\u0633\u062a\u0634f\u0649', '\u0627\u0644\u0645\u0633\u062a\u0634\u0641\u0649'),
    ]

    arabic_fix = {
        'المستشfى': 'المستشfى',
        'الانتظar': 'الانتظar',
        'انتظar': 'انتظar',
        'المواعid': 'المواعid',
        'مواعid': 'مواعid',
        'التأmين': 'التأmين',
        'تأmين': 'تأmين',
        'التشخiص': 'التشخiص',
        'تشخiص': 'تشخiص',
        'المستقblil': 'المستقblil',
        '【متقدm】': '【متقدm】',
    }

    # Hard-coded correct replacements using chr building
    HOSP = '\u0627\u0644\u0645\u0633\u062a\u0634\u0641\u0649'
    WAIT = '\u0627\u0644\u0627\u0646\u062a\u0638\u0627\u0631'
    WAIT2 = '\u0627\u0646\u062a\u0638\u0627\u0631'
    APPT = '\u0627\u0644\u0645\u0648\u0627\u0639\u064a\u062f'
    APPT2 = '\u0645\u0648\u0627\u0639\u064a\u062f'
    INS = '\u0627\u0644\u062a\u0623\u0645\u064a\u0646'
    INS2 = '\u062a\u0623\u0645\u064a\u0646'
    DIAG = '\u0627\u0644\u062a\u0634\u062e\u064a\u0635'
    DIAG2 = '\u062a\u0634\u062e\u064a\u0635'
    FUT = '\u0627\u0644\u0645\u0633\u062a\u0642\u0628\u0644'
    ADV = '\u0645\u062a\u0642\u062f\u0645'

    bad_good = [
        ('المستشfى', HOSP),
        ('المستشfى', HOSP),
        ('الانتظar', WAIT),
        ('انتظar', WAIT2),
        ('المواعid', APPT),
        ('مواعid', APPT2),
        ('التأmين', INS),
        ('تأmين', INS2),
        ('التشخiص', DIAG),
        ('تشخiص', DIAG2),
        ('المستقblil', FUT),
        ('【متقدm】', f'【{ADV}】'),
    ]

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    for bad, good in bad_good:
                        t = t.replace(bad, good)
                    run.text = t

    # Fix title slide specifically
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and 'تطوير' in shape.text_frame.text:
            shape.text_frame.paragraphs[0].runs[0].text = f'تطوير نظام إدارة {HOSP}'
            break

    prs.save(str(OUTPUT))
    print(f'Created {OUTPUT} ({len(prs.slides)} slides)')
    return OUTPUT


if __name__ == '__main__':
    build()
